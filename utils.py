
import os, io, time, json, math, base64, hashlib
from typing import List, Tuple, Dict
import numpy as np
import pandas as pd
from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions
from datetime import datetime, timedelta
import requests

ACCOUNT = os.getenv("AZURE_STORAGE_ACCOUNT")
OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT").rstrip("/")
EMBED_DEPLOYMENT = os.getenv("AZURE_OPENAI_EMBED_DEPLOYMENT")
CHAT_DEPLOYMENT = os.getenv("AZURE_OPENAI_CHAT_DEPLOYMENT")
SUPPORT_URL = os.getenv("SUPPORT_TICKET_URL", "https://support.example.com/new")

BS = BlobServiceClient(account_url=f"https://{ACCOUNT}.blob.core.windows.net", credential=DefaultAzureCredential())

def get_text_from_pdf(stream: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(stream))
    texts = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            texts.append("")
    return "\n\n".join(texts)

def get_text_from_docx(stream: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(stream))
    return "\n".join([p.text for p in doc.paragraphs])

def get_text_from_pptx(stream: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(stream))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def smart_chunk(text: str, max_chars: int = 2500, overlap: int = 250) -> List[str]:
    text = text.replace("\r", "")
    chunks = []
    i = 0
    while i < len(text):
        j = min(i + max_chars, len(text))
        # try to break on a paragraph boundary
        k = text.rfind("\n\n", i, j)
        if k == -1 or j - k < 400:
            k = j
        chunks.append(text[i:k].strip())
        i = max(k - overlap, k)
    return [c for c in chunks if c]

def _aad_token() -> str:
    cred = DefaultAzureCredential()
    token = cred.get_token("https://cognitiveservices.azure.com/.default")
    return token.token

def embed_texts(texts: List[str]) -> np.ndarray:
    # Azure OpenAI embeddings via AAD bearer
    url = f"{OPENAI_ENDPOINT}/openai/deployments/{EMBED_DEPLOYMENT}/embeddings?api-version=2024-02-15-preview"
    headers = {"Authorization": f"Bearer {_aad_token()}","Content-Type":"application/json"}
    vecs = []
    for i in range(0, len(texts), 32):
        batch = texts[i:i+32]
        payload = {"input": batch}
        r = requests.post(url, headers=headers, json=payload, timeout=60)
        r.raise_for_status()
        data = r.json()
        for item in data["data"]:
            vecs.append(np.array(item["embedding"], dtype=np.float32))
    arr = np.vstack(vecs)
    # normalize for cosine
    norms = np.linalg.norm(arr, axis=1, keepdims=True) + 1e-8
    arr = arr / norms
    return arr

def save_parquet(container: str, name: str, df: pd.DataFrame):
    bs = BS.get_container_client(container)
    buf = io.BytesIO()
    df.to_parquet(buf, index=False)
    buf.seek(0)
    bs.upload_blob(name, buf, overwrite=True)

def load_all_embeddings() -> Tuple[pd.DataFrame, np.ndarray]:
    cc = BS.get_container_client("sop-index")
    # Expect two blobs: chunks.parquet, embs.npy
    chunks = pd.read_parquet(io.BytesIO(cc.download_blob("chunks.parquet").readall()))
    embs = np.load(io.BytesIO(cc.download_blob("embs.npy").readall()))
    return chunks, embs

def build_and_store_faiss(df: pd.DataFrame, embs: np.ndarray):
    import faiss
    d = embs.shape[1]
    nlist = 256
    m = 32
    quantizer = faiss.IndexFlatIP(d)
    index = faiss.IndexIVFPQ(quantizer, d, nlist, m, 8)
    index.train(embs)
    index.add(embs)
    # persist
    cc = BS.get_container_client("sop-index")
    # index
    idx_bytes = faiss.serialize_index(index)
    cc.upload_blob("index.faiss", idx_bytes, overwrite=True)
    # mapping
    save_parquet("sop-index", "chunks.parquet", df)
    # raw embs (for future rebuilds)
    buf = io.BytesIO()
    np.save(buf, embs)
    buf.seek(0)
    cc.upload_blob("embs.npy", buf.getvalue(), overwrite=True)
    # stats
    stats = {"dim": int(d), "nlist": int(nlist), "m": int(m), "count": int(embs.shape[0]), "savedAt": datetime.utcnow().isoformat()+"Z"}
    cc.upload_blob("stats.json", json.dumps(stats), overwrite=True)

def faiss_search(qvec: np.ndarray, topk: int = 8) -> List[int]:
    import faiss
    cc = BS.get_container_client("sop-index")
    idx_bytes = cc.download_blob("index.faiss").readall()
    index = faiss.deserialize_index(idx_bytes)
    D, I = index.search(qvec.astype("float32"), topk)
    return I[0].tolist()

def make_sas(container: str, name: str, minutes: int = 20) -> str:
    # user delegation SAS (works with AAD/MI)
    key = BS.get_user_delegation_key(datetime.utcnow(), datetime.utcnow() + timedelta(hours=1))
    sas = generate_blob_sas(
        account_name=ACCOUNT,
        container_name=container,
        blob_name=name,
        permission=BlobSasPermissions(read=True),
        expiry=datetime.utcnow() + timedelta(minutes=minutes),
        user_delegation_key=key
    )
    return f"https://{ACCOUNT}.blob.core.windows.net/{container}/{name}?{sas}"

def rate_limit_ok(conversation_id: str) -> Tuple[bool, int]:
    bc = BS.get_container_client("bot-state")
    blob = f"{conversation_id}.json"
    try:
        data = json.loads(bc.download_blob(blob).readall())
    except Exception:
        data = {"count": 0, "last": time.time()}
    # reset after 2 hours
    if time.time() - data.get("last", 0) > 7200:
        data = {"count": 0, "last": time.time()}
    if data["count"] >= 3:
        return False, data["count"]
    data["count"] += 1
    data["last"] = time.time()
    bc.upload_blob(blob, json.dumps(data), overwrite=True)
    return True, data["count"]

def azure_openai_answer(question: str, context: str) -> str:
    url = f"{OPENAI_ENDPOINT}/openai/deployments/{CHAT_DEPLOYMENT}/chat/completions?api-version=2024-02-15-preview"
    headers = {"Authorization": f"Bearer {_aad_token()}","Content-Type":"application/json"}
    sys = "You are an SOP assistant. Answer ONLY using the provided context. If insufficient, say you cannot find it and suggest opening a support ticket."
    payload = {"messages":[{"role":"system","content":sys},{"role":"user","content":f"QUESTION:\n{question}\n\nCONTEXT:\n{context}"}],
               "temperature":0.2}
    r = requests.post(url, headers=headers, json=payload, timeout=60)
    r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"]
